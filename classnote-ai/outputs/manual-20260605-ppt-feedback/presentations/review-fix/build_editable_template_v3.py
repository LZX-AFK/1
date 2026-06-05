from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(r"D:\classnote-ai-hx\outputs\manual-20260605-ppt-feedback\presentations\review-fix")
OUT_DIR = ROOT / "output"
MEDIA = ROOT / "pptx-unpack-source" / "ppt" / "media"
ASSETS = ROOT / "assets" / "source-slide7"
OUT = OUT_DIR / "超硅氨水项目复盘报告_模板视觉三页可编辑.pptx"
OUT_DIR.mkdir(parents=True, exist_ok=True)

W, H = 13.333, 7.5
BLUE = RGBColor(0, 65, 113)
GREEN = RGBColor(101, 190, 49)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(93, 101, 112)
BORDER = RGBColor(189, 211, 228)
LIGHT = RGBColor(244, 249, 252)
HEADER = RGBColor(0, 68, 125)
WHITE = RGBColor(255, 255, 255)
CREAM = RGBColor(255, 249, 236)
GOLD = RGBColor(228, 188, 108)


def add_text(slide, text, x, y, w, h, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    tf.vertical_anchor = valign
    paragraphs = str(text).split("\n")
    for idx, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = para
        run.font.name = "Microsoft YaHei"
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def rect(slide, x, y, w, h, fill=WHITE, line=BORDER, radius=False, weight=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(weight)
    return shp


def oval(slide, x, y, w, h, fill=GREEN, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    return shp


def label(slide, text, x, y, w=1.02, h=0.35, fill=GREEN):
    shp = rect(slide, x, y, w, h, fill=fill, line=None, radius=True)
    tf = shp.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(12)
    r.font.bold = True
    r.font.color.rgb = WHITE
    return shp


def add_picture(slide, path, x, y, w=None, h=None):
    if not Path(path).exists():
        return None
    kwargs = {}
    if w is not None:
        kwargs["width"] = Inches(w)
    if h is not None:
        kwargs["height"] = Inches(h)
    return slide.shapes.add_picture(str(path), Inches(x), Inches(y), **kwargs)


def base(slide, title, page, subtitle=None):
    rect(slide, 0, 0.40, 0.34, 0.46, fill=BLUE, line=None)
    rect(slide, 0.43, 0.40, 0.15, 0.46, fill=BLUE, line=None)
    add_text(slide, title, 0.74, 0.50, 7.4, 0.36, 20, TEXT, True)
    if subtitle:
        add_text(slide, subtitle, 0.77, 0.84, 8.6, 0.20, 8.5, MUTED)
    add_picture(slide, MEDIA / "image14.png", 11.08, 0.32, w=1.25)
    add_text(slide, "版权所有 © 2025 Capchem. All Rights Reserved.", 0.46, 7.17, 3.2, 0.14, 8.5, TEXT)
    rect(slide, 3.42, 7.24, 8.80, 0.01, fill=RGBColor(125, 125, 125), line=None)
    add_text(slide, str(page), 12.28, 7.04, 0.16, 0.16, 8.5, MUTED, align=PP_ALIGN.RIGHT)


def slide_problem_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base(s, "四、问题与整改措施", 6, "在原表格框架上调整口径：弱化个人批判，强化流程闭环与商务协同。")
    x0, y0, w = 0.75, 1.25, 9.84
    row_h, header_h = 0.90, 0.45
    col = [0.75, 2.92, 5.67, 10.59]
    rect(s, x0, y0, w, header_h, fill=HEADER, line=None)
    for x in col[1:-1]:
        rect(s, x, y0, 0.01, header_h + row_h * 4, fill=BORDER, line=None)
    for i, htxt in enumerate(["问题点", "调整后原因", "整改动作"]):
        add_text(s, htxt, col[i] + 0.12, y0 + 0.13, col[i + 1] - col[i] - 0.24, 0.22, 12, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    rows = [
        ("小量测试闭环不足", "测试结束后，对空桶回收风险节点识别不充分，未及时形成回收确认。", "测试完成后3个工作日内核对空桶状态，建立个人跟进台账。"),
        ("中量测试跟进滞后", "沿用以往测试经验，对本次空桶回收特殊性预判不足，沟通节点偏晚。", "发货前完成客户采购、厂务、仓库三方邮件确权，锁定回收窗口。"),
        ("商务信息同步不足", "商务条款与回收要求未形成统一书面同步，存在跨部门信息断点。", "销售与商务共同确认关键事项，客户反馈5个工作日内闭环，逾期升级。"),
        ("异常响应不够及时", "发现空桶进入处置流程后，缺少快速联动机制和责任边界。", "建立异常升级机制，联动商务、客户仓库及相关部门及时处置。"),
    ]
    for r, row in enumerate(rows):
        y = y0 + header_h + r * row_h
        fill = RGBColor(249, 252, 254) if r % 2 else WHITE
        rect(s, x0, y, w, row_h, fill=fill, line=BORDER, weight=0.5)
        for x in col[1:-1]:
            rect(s, x, y, 0.01, row_h, fill=BORDER, line=None)
        oval(s, 0.90, y + 0.29, 0.30, 0.30, fill=GREEN)
        add_text(s, str(r + 1), 0.98, y + 0.33, 0.12, 0.12, 10, WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, row[0], 1.45, y + 0.31, 1.48, 0.34, 12.5, TEXT, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        add_text(s, row[1], 3.08, y + 0.18, 2.45, 0.63, 10.5, TEXT)
        add_text(s, row[2], 5.82, y + 0.18, 4.45, 0.63, 10.5, TEXT)
    rect(s, 10.84, 1.26, 1.63, 3.62, fill=LIGHT, line=BORDER, radius=True)
    add_picture(s, MEDIA / "image7.png", 11.25, 1.52, w=0.72)
    add_text(s, "事实总结", 11.20, 2.57, 0.88, 0.24, 14, BLUE, True, PP_ALIGN.CENTER)
    add_text(s, "44个空桶报废\n8万元直接损失\n整改重点转向流程闭环与部门协同", 11.03, 2.95, 1.25, 1.02, 12, TEXT, align=PP_ALIGN.CENTER)
    rect(s, 0.75, 5.43, 11.70, 0.77, fill=LIGHT, line=BORDER, radius=True)
    label(s, "核心调整", 1.00, 5.67, w=1.00, h=0.36, fill=BLUE)
    add_text(s, "不再强化“个人严重失职”表述，改为“风险识别不足、节点预判不足、协同机制不足”，并以流程机制补强作为整改主线。", 2.25, 5.62, 9.65, 0.44, 15, TEXT)


def slide_evidence(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base(s, "客户反馈补充证据与整改闭环", 7, "保留原三项证据，但按模板正文页重排为“证据-判断-整改结论”。")
    rect(s, 0.65, 1.27, 4.43, 3.90, fill=WHITE, line=BORDER, radius=True)
    label(s, "证据 01", 0.90, 1.48, w=1.03, h=0.35)
    add_text(s, "4月22日：报价阶段已明确风险条款", 0.90, 2.14, 3.55, 0.32, 17, TEXT, True)
    add_text(s, "报价单备注已提示包装桶妥善保管及丢失赔偿要求，说明风险条款存在，但未进入后续发货、测试、回收执行闭环。", 0.90, 2.52, 3.70, 0.68, 12, MUTED)
    rect(s, 1.10, 3.38, 3.55, 1.25, fill=RGBColor(248, 248, 248), line=RGBColor(226, 226, 226), radius=True)
    add_picture(s, ASSETS / "shape-8.png", 1.42, 3.50, w=2.95)
    cards = [
        (5.50, 1.27, "证据 02", "5月18日：首次沟通空桶回收安排", "客户反馈“送货时带回”，说明沟通确实发生，但节点已偏晚，未形成发货前锁定与回收确认。", "shape-14.png"),
        (5.50, 3.23, "证据 03", "5月21日：确认已处理且无法追回", "客户反馈仓库已处理、空桶追不回，最终形成44个空桶报废及8万元直接损失。", "shape-18.png"),
    ]
    for x, y, tag, title, body, pic in cards:
        rect(s, x, y, 6.91, 1.67, fill=WHITE, line=BORDER, radius=True)
        label(s, tag, x + 0.23, y + 0.22, w=1.02, h=0.35)
        add_text(s, title, x + 1.50, y + 0.28, 3.90, 0.32, 17, TEXT, True)
        add_text(s, body, x + 1.50, y + 0.65, 3.82, 0.54, 11.5, MUTED)
        rect(s, x + 5.42, y + 0.45, 1.18, 0.73, fill=RGBColor(248, 248, 248), line=RGBColor(226, 226, 226), radius=True)
        add_picture(s, ASSETS / pic, x + 5.50, y + 0.52, w=1.03)
    rect(s, 0.65, 5.50, 11.76, 0.85, fill=LIGHT, line=BORDER, radius=True)
    label(s, "复盘结论", 0.92, 5.75, w=1.06, h=0.36, fill=BLUE)
    add_text(s, "本次问题不是“没有沟通”，而是风险条款未进入执行闭环，且沟通节点明显滞后。整改重点由“事后询问”前移为“发货前确权、测试中跟踪、测试后限时回收”。", 2.23, 5.67, 9.82, 0.48, 14, TEXT)


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    base(s, "五、复盘总结", 8, "保留原“后续计划 + 复盘反思”框架，融入商务协同与损失挽回计划。")
    rect(s, 0.65, 1.25, 2.23, 3.92, fill=BLUE, line=None)
    add_picture(s, MEDIA / "image20.png", 0.79, 1.35, w=1.95)
    add_picture(s, MEDIA / "image21.png", 0.79, 2.67, w=1.95)
    rect(s, 0.79, 4.20, 1.95, 0.87, fill=WHITE, line=None)
    add_text(s, "复盘反思 · 明确问题\n总结改进 · 持续优化", 1.02, 4.42, 1.50, 0.45, 13, BLUE, True, PP_ALIGN.CENTER)
    blocks = [
        ("后续工作计划", "1. 三季度大批量测试前，完成客户采购、厂务、仓库三方对接及邮件确权。\n2. 发货前锁定空桶暂存、回收窗口、赔付条款和客户对接人。\n3. 测试期间建立节点台账，测试后3个工作日内核对空桶状态。\n4. 关键事项发出后5个工作日内需取得客户回复，逾期联动商务升级。", "image10.png", 1.28),
        ("复盘反思", "本次问题暴露出本人对测试空桶回收风险识别不充分，对客户回收节点和内部协同机制预判不足，导致跟进节奏滞后。因以往测试较少涉及空桶回收要求，本人沿用既有测试经验，对本次特殊风险节点识别不足。", "image12.png", 1.18),
        ("商务协同与损失挽回", "后续由销售牵头、商务协同，将回收要求和客户反馈时限纳入工作机制。结合三季度大批量测试推进，通过价格策略优化、商务条款调整等方式，争取逐步抵回本次8万元空桶损失。", "image9.png", 1.18),
    ]
    y = 1.25
    for title, body, icon, h in blocks:
        rect(s, 3.17, y, 9.24, h, fill=WHITE, line=BORDER, radius=True)
        rect(s, 3.38, y + 0.28, 0.52, 0.52, fill=LIGHT, line=BORDER, radius=True)
        add_picture(s, MEDIA / icon, 3.48, y + 0.37, w=0.32)
        add_text(s, title, 4.17, y + 0.20, 3.10, 0.30, 16, BLUE, True)
        add_text(s, body, 4.17, y + 0.52, 8.05, h - 0.62, 10.7 if h > 1.2 else 11.3, TEXT)
        y += h + 0.22
    rect(s, 3.17, 5.38, 9.24, 0.86, fill=CREAM, line=GOLD, radius=True)
    label(s, "表述口径", 3.42, 5.65, w=1.10, h=0.36, fill=BLUE)
    add_text(s, "承认识别不足与跟进滞后，但避免过度个人化；整改落点放在流程闭环、部门协同、节点审核和商务挽回。", 4.75, 5.58, 7.25, 0.42, 14.2, TEXT)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
slide_problem_table(prs)
slide_evidence(prs)
slide_summary(prs)
prs.save(OUT)
print(OUT)
