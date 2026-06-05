from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(r"D:\classnote-ai-hx\outputs\manual-20260605-ppt-feedback\presentations\review-fix")
OUT_DIR = ROOT / "output"
PREVIEW_DIR = ROOT / "preview" / "final"
ASSETS = ROOT / "assets" / "source-slide7"
LOGO = ROOT / "pptx-unpack-template" / "ppt" / "media" / "image2.png"
OUT = OUT_DIR / "超硅氨水项目复盘报告_反馈整改三页.pptx"

OUT_DIR.mkdir(parents=True, exist_ok=True)
PREVIEW_DIR.mkdir(parents=True, exist_ok=True)

W, H = 13.333, 7.5
PX_W, PX_H = 1600, 900

BLUE = RGBColor(0, 64, 112)
DEEP = RGBColor(0, 43, 85)
GREEN = RGBColor(0, 176, 80)
LIGHT = RGBColor(244, 248, 251)
LINE = RGBColor(199, 214, 225)
GRAY = RGBColor(88, 96, 105)
BLACK = RGBColor(31, 41, 55)


def rgb_tuple(c):
    return c[0], c[1], c[2]


def add_text(slide, text, x, y, w, h, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    box.text_frame.clear()
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    p = box.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_rect(slide, x, y, w, h, fill, line=None, radius=True, weight=1):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(weight)
    return shp


def add_label(slide, text, x, y):
    tag = add_rect(slide, x, y, 0.92, 0.34, GREEN, None, True)
    tag.text_frame.clear()
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tag.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    return tag


def add_base(slide, title, subtitle, page_no):
    add_rect(slide, 0.0, 0.39, 0.34, 0.45, BLUE, None, False)
    add_rect(slide, 0.43, 0.39, 0.16, 0.45, BLUE, None, False)
    add_text(slide, title, 0.68, 0.43, 6.8, 0.36, 19, BLACK, True)
    add_text(slide, subtitle, 0.86, 0.78, 8.5, 0.22, 9.5, BLACK)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.0), Inches(0.23), width=Inches(1.38))
    add_text(slide, "版权所有 © 2025 Capchem. All Rights Reserved.", 0.47, 7.14, 3.25, 0.16, 8.5, BLACK)
    line = add_rect(slide, 3.46, 7.19, 8.72, 0.012, RGBColor(125, 125, 125), None, False)
    add_text(slide, str(page_no), 12.12, 7.03, 0.18, 0.14, 8, GRAY, False, PP_ALIGN.RIGHT)


def add_circle_text(slide, text, x, y, size, fill, font_size=14):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.fill.background()
    if text:
        c.text_frame.clear()
        c.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(font_size)
        r.font.bold = True
        r.font.color.rgb = RGBColor(255, 255, 255)
    return c


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# Slide 1
s = prs.slides.add_slide(blank)
add_base(
    s,
    "客户反馈补充证据与整改闭环",
    "补充报价条款、首次沟通与客户最终反馈，明确问题从“沟通是否发生”转为“回收闭环是否前置锁定”。",
    7,
)
add_rect(s, 0.78, 1.30, 12.0, 4.35, RGBColor(255, 255, 255), LINE, True)
add_rect(s, 1.20, 2.18, 10.30, 0.025, BLUE, None, False)
for x, date, phase, claim in [
    (1.76, "4月22日", "报价阶段", "风险条款已提示"),
    (5.28, "5月18日", "首次沟通", "回收安排节点偏晚"),
    (8.86, "5月21日", "客户确认", "空桶已处理，无法追回"),
]:
    add_text(s, date, x - 0.45, 1.50, 1.2, 0.25, 18, BLACK, True, PP_ALIGN.CENTER)
    add_text(s, phase, x - 0.42, 1.82, 1.15, 0.22, 10.5, GRAY, False, PP_ALIGN.CENTER)
    add_circle_text(s, "", x, 2.02, 0.38, GREEN, 14)
    add_text(s, claim, x - 0.55, 2.50, 1.7, 0.35, 11.5, BLACK, True, PP_ALIGN.CENTER)
for x, img in zip([1.23, 4.72, 8.25], ["shape-8.png", "shape-14.png", "shape-18.png"]):
    add_rect(s, x, 3.10, 3.10, 0.88, RGBColor(248, 248, 248), RGBColor(231, 231, 231), True)
    p = ASSETS / img
    if p.exists():
        s.shapes.add_picture(str(p), Inches(x + 0.72), Inches(3.25), width=Inches(1.65))
add_label(s, "复盘结论", 1.16, 4.76)
add_text(
    s,
    "本次问题不是“没有沟通”，而是报价风险条款未进入执行闭环，且沟通节点明显滞后。整改重点前移至“发货前确权、测试中跟踪、测试后限时回收”。",
    2.25,
    4.74,
    9.55,
    0.45,
    13.5,
    BLACK,
)

# Slide 2
s = prs.slides.add_slide(blank)
add_base(
    s,
    "复盘总结：个人识别不足与机制补强",
    "按老板反馈弱化个人批判，聚焦风险识别、经验不足与跨部门协同机制。",
    8,
)
add_rect(s, 0.82, 1.40, 5.62, 4.42, RGBColor(255, 255, 255), LINE, True)
add_rect(s, 7.07, 1.40, 5.62, 4.42, RGBColor(255, 255, 255), LINE, True)
add_label(s, "问题表述", 1.18, 1.73)
add_text(
    s,
    "本次项目造成 44 个空桶报废及 8 万元直接损失，暴露出本人对测试空桶回收风险识别不充分，对客户回收节点和内部协同机制预判不足，导致跟进节奏滞后。",
    1.18,
    2.30,
    4.65,
    0.90,
    15,
    BLACK,
)
add_text(
    s,
    "因以往测试流程中较少涉及空桶回收要求，本人沿用既有测试跟进经验，对本次空桶回收的特殊性和关键风险节点识别不足。",
    1.18,
    3.70,
    4.65,
    0.85,
    15,
    BLACK,
)
add_label(s, "改进方向", 7.45, 1.73)
for idx, txt in enumerate(
    [
        "前置识别：发货前确认空桶归属、暂存、回收时限及赔付条款。",
        "过程跟踪：测试期间建立节点台账，避免仅依赖单点转述。",
        "协同补强：销售、商务、厂务/仓库信息同步，关键事项书面留痕。",
        "审核空间：重大风险节点提交部门复核，减少个人经验判断偏差。",
    ]
):
    y = 2.25 + idx * 0.72
    add_circle_text(s, "", 7.45, y + 0.05, 0.17, GREEN, 1)
    add_text(s, txt, 7.76, y, 4.45, 0.48, 13.5, BLACK)
add_rect(s, 1.16, 5.37, 10.98, 0.018, BLUE, None, False)
add_text(
    s,
    "表述基调：承认识别不足与跟进滞后，但将整改落点放在流程闭环、部门协同与节点审核。",
    2.75,
    6.07,
    8.7,
    0.25,
    13.5,
    BLACK,
    True,
    PP_ALIGN.CENTER,
)

# Slide 3
s = prs.slides.add_slide(blank)
add_base(
    s,
    "后续工作计划：商务协同与损失挽回",
    "把“损失挽回计划”融入正式工作机制，避免底部硬塞式新增。",
    9,
)
steps = [
    ("01", "发货前确权", "销售牵头，商务协同客户确认空桶回收、赔付条款、回收窗口及对接人。"),
    ("02", "5个工作日反馈", "关键事项发出后 5 个工作日内需取得客户回复；未回复即升级联动跟进。"),
    ("03", "台账与预警", "建立测试物料与空桶回收台账，测试完成后 3 个工作日内核对空桶状态。"),
    ("04", "损失挽回计划", "结合三季度大批量测试，通过价格策略和商务条款，争取抵回 8 万元损失。"),
]
for idx, (num, head, body) in enumerate(steps):
    x = 0.98 + idx * 3.02
    add_rect(s, x, 1.65, 2.65, 3.40, RGBColor(255, 255, 255), LINE, True)
    add_circle_text(s, num, x + 0.38, 1.97, 0.58, GREEN, 16)
    add_text(s, head, x + 0.35, 2.85, 1.95, 0.30, 15.5, BLACK, True, PP_ALIGN.CENTER)
    add_text(s, body, x + 0.35, 3.40, 1.95, 1.05, 12.0, BLACK, False, PP_ALIGN.CENTER)
add_rect(s, 2.20, 5.35, 8.98, 0.02, BLUE, None, False)
for x, label in [(2.18, "发货前"), (5.05, "测试中"), (7.93, "测试后"), (10.82, "商务回收")]:
    add_circle_text(s, "", x, 5.20, 0.32, BLUE, 1)
    add_text(s, label, x - 0.34, 5.67, 1.0, 0.18, 10.5, GRAY, False, PP_ALIGN.CENTER)
add_text(
    s,
    "原则：风险条款前置、客户反馈限时、逾期协同升级、损失挽回纳入后续商务方案。",
    3.12,
    6.22,
    7.4,
    0.24,
    13.5,
    BLACK,
    True,
    PP_ALIGN.CENTER,
)

prs.save(OUT)


def load_font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
    ]
    for cand in candidates:
        if Path(cand).exists():
            return ImageFont.truetype(cand, size)
    return ImageFont.load_default()


def draw_wrapped(draw, text, box, font, fill, align="left", line_spacing=1.18):
    x, y, w, h = box
    lines = []
    current = ""
    for ch in text:
        trial = current + ch
        if draw.textlength(trial, font=font) <= w or not current:
            current = trial
        else:
            lines.append(current)
            current = ch
    if current:
        lines.append(current)
    line_h = int(font.size * line_spacing)
    for i, line in enumerate(lines):
        yy = y + i * line_h
        if yy > y + h - line_h:
            break
        xx = x
        if align == "center":
            xx = x + (w - draw.textlength(line, font=font)) / 2
        draw.text((xx, yy), line, font=font, fill=fill)


def preview_base(draw, title, subtitle, page_no):
    draw.rectangle([0, 47, 40, 102], fill=rgb_tuple(BLUE))
    draw.rectangle([52, 47, 70, 102], fill=rgb_tuple(BLUE))
    draw.text((82, 52), title, font=load_font(32, True), fill=rgb_tuple(BLACK))
    draw.text((103, 94), subtitle, font=load_font(15), fill=rgb_tuple(BLACK))
    if LOGO.exists():
        logo = Image.open(LOGO).convert("RGBA")
        logo.thumbnail((172, 52))
        img.paste(logo, (1320, 28), logo)
    draw.text((56, 861), "版权所有 © 2025 Capchem. All Rights Reserved.", font=load_font(17), fill=rgb_tuple(BLACK))
    draw.line([410, 869, 1460, 869], fill=(125, 125, 125), width=2)
    draw.text((1460, 845), str(page_no), font=load_font(15), fill=rgb_tuple(GRAY))


for idx in range(1, 4):
    img = Image.new("RGB", (PX_W, PX_H), "white")
    d = ImageDraw.Draw(img)
    if idx == 1:
        preview_base(d, "客户反馈补充证据与整改闭环", "补充报价条款、首次沟通与客户最终反馈，明确问题从“沟通是否发生”转为“回收闭环是否前置锁定”。", 7)
        d.rounded_rectangle([94, 156, 1534, 678], radius=70, outline=rgb_tuple(LINE), width=2)
        d.rectangle([144, 262, 1380, 265], fill=rgb_tuple(BLUE))
        for x, date, phase, claim in [(264, "4月22日", "报价阶段", "风险条款已提示"), (634, "5月18日", "首次沟通", "回收安排节点偏晚"), (1064, "5月21日", "客户确认", "空桶已处理，无法追回")]:
            d.text((x - 68, 180), date, font=load_font(34, True), fill=rgb_tuple(BLACK))
            d.text((x - 42, 220), phase, font=load_font(18), fill=rgb_tuple(GRAY))
            d.ellipse([x, 242, x + 46, 288], fill=rgb_tuple(GREEN))
            draw_wrapped(d, claim, (x - 66, 300, 190, 48), load_font(20, True), rgb_tuple(BLACK), "center")
        for x, name in zip([147, 566, 990], ["shape-8.png", "shape-14.png", "shape-18.png"]):
            d.rounded_rectangle([x, 372, x + 372, 478], radius=28, fill=(248, 248, 248), outline=(231, 231, 231), width=2)
            p = ASSETS / name
            if p.exists():
                ev = Image.open(p).convert("RGBA")
                ev.thumbnail((198, 80))
                img.paste(ev, (x + 88, 389), ev)
        d.rounded_rectangle([139, 571, 250, 610], radius=10, fill=rgb_tuple(GREEN))
        d.text((162, 582), "复盘结论", font=load_font(18, True), fill="white")
        draw_wrapped(d, "本次问题不是“没有沟通”，而是报价风险条款未进入执行闭环，且沟通节点明显滞后。整改重点前移至“发货前确权、测试中跟踪、测试后限时回收”。", (270, 570, 1140, 58), load_font(25), rgb_tuple(BLACK))
    elif idx == 2:
        preview_base(d, "复盘总结：个人识别不足与机制补强", "按老板反馈弱化个人批判，聚焦风险识别、经验不足与跨部门协同机制。", 8)
        d.rounded_rectangle([98, 168, 773, 698], radius=70, outline=rgb_tuple(LINE), width=2)
        d.rounded_rectangle([848, 168, 1523, 698], radius=70, outline=rgb_tuple(LINE), width=2)
        d.rounded_rectangle([142, 208, 252, 247], radius=9, fill=rgb_tuple(GREEN))
        d.text((162, 219), "问题表述", font=load_font(18, True), fill="white")
        draw_wrapped(d, "本次项目造成 44 个空桶报废及 8 万元直接损失，暴露出本人对测试空桶回收风险识别不充分，对客户回收节点和内部协同机制预判不足，导致跟进节奏滞后。", (142, 276, 560, 110), load_font(28), rgb_tuple(BLACK))
        draw_wrapped(d, "因以往测试流程中较少涉及空桶回收要求，本人沿用既有测试跟进经验，对本次空桶回收的特殊性和关键风险节点识别不足。", (142, 444, 560, 110), load_font(28), rgb_tuple(BLACK))
        d.rounded_rectangle([894, 208, 1004, 247], radius=9, fill=rgb_tuple(GREEN))
        d.text((914, 219), "改进方向", font=load_font(18, True), fill="white")
        for i, txt in enumerate(["前置识别：发货前确认空桶归属、暂存、回收时限及赔付条款。", "过程跟踪：测试期间建立节点台账，避免仅依赖单点转述。", "协同补强：销售、商务、厂务/仓库信息同步，关键事项书面留痕。", "审核空间：重大风险节点提交部门复核，减少个人经验判断偏差。"]):
            y = 270 + i * 86
            d.ellipse([894, y + 8, 916, y + 30], fill=rgb_tuple(GREEN))
            draw_wrapped(d, txt, (932, y, 535, 58), load_font(25), rgb_tuple(BLACK))
        d.rectangle([139, 644, 1458, 647], fill=rgb_tuple(BLUE))
        d.text((330, 728), "表述基调：承认识别不足与跟进滞后，但将整改落点放在流程闭环、部门协同与节点审核。", font=load_font(27, True), fill=rgb_tuple(BLACK))
    else:
        preview_base(d, "后续工作计划：商务协同与损失挽回", "把“损失挽回计划”融入正式工作机制，避免底部硬塞式新增。", 9)
        for i, (num, head, body) in enumerate(steps):
            x = 118 + i * 362
            d.rounded_rectangle([x, 198, x + 318, 606], radius=50, outline=rgb_tuple(LINE), width=2)
            d.ellipse([x + 46, 236, x + 116, 306], fill=rgb_tuple(GREEN))
            d.text((x + 63, 255), num, font=load_font(30, True), fill="white")
            d.text((x + 84, 341), head, font=load_font(28, True), fill=rgb_tuple(BLACK))
            draw_wrapped(d, body, (x + 42, 408, 234, 130), load_font(24), rgb_tuple(BLACK), "center")
        d.rectangle([264, 642, 1342, 645], fill=rgb_tuple(BLUE))
        for x, label in [(262, "发货前"), (606, "测试中"), (952, "测试后"), (1298, "商务回收")]:
            d.ellipse([x, 624, x + 38, 662], fill=rgb_tuple(BLUE))
            d.text((x - 28, 682), label, font=load_font(20), fill=rgb_tuple(GRAY))
        d.text((374, 746), "原则：风险条款前置、客户反馈限时、逾期协同升级、损失挽回纳入后续商务方案。", font=load_font(27, True), fill=rgb_tuple(BLACK))
    img.save(PREVIEW_DIR / f"fixed-slide-{idx:02d}.png")

print(OUT)
print(PREVIEW_DIR)
